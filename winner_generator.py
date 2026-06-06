# -*- coding: utf-8 -*-
"""
윈도우용 경품 당첨자 PPT/CSV 생성 프로그램

필요 파일(같은 폴더):
- 전체.csv
- 블락.csv
- 상품.csv
- 당첨자.pptx

출력:
- ./당첨자/당첨자_N.pptx
- ./블락/블락_N.csv
"""

import csv
import os
import random
import re
import sys
import traceback
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt
except Exception as e:
    Presentation = None

try:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk
except Exception:
    tk = None

# 필요하면 여기만 바꾸면 됩니다.
OUTPUT_NAME_FORMAT = "{name}({last4})"   # 예: 홍길동(1234)
WINNER_SEPARATOR = " "                   # 당첨자 사이 구분자: 띄어쓰기 1번
MAX_WINNERS_PER_LINE = 5                 # PPT 한 줄당 최대 당첨자 수
MAX_WINNERS_PER_SLIDE = 25                # PPT 한 장당 최대 당첨자 수
FORCE_NO_DUPLICATE = True                 # 전체 상품 기준 중복 당첨 항상 방지


NAME_HEADERS = ["이름", "성명", "신청자", "참가자", "회원명", "당첨자명", "name"]
PHONE_HEADERS = ["휴대전화", "휴대폰", "핸드폰", "전화번호", "연락처", "휴대전화번호", "phone", "mobile", "tel"]
PRODUCT_HEADERS = ["상품명", "상품", "경품", "product", "prize"]
COUNT_HEADERS = ["인원수", "인원", "수량", "당첨인원", "count", "qty", "quantity"]
DUP_HEADERS = ["중복허용", "중복 허용", "중복", "allow_duplicate", "duplicate"]

YES_VALUES = {"y", "yes", "true", "1", "가능", "허용", "중복허용", "o", "○", "예", "네"}
NO_VALUES = {"n", "no", "false", "0", "불가", "불허", "중복불가", "x", "×", "아니오", "아니요"}


@dataclass
class Person:
    source_row_no: int  # 전체.csv에서 실제 몇 번째 줄인지. 헤더가 1줄이면 첫 데이터는 2.
    name: str
    phone: str
    raw_row: Dict[str, str] = None  # 전체.csv 원본 행. 블락_N.csv 생성 시 사용

    @property
    def key(self) -> Tuple[str, str]:
        return normalize_name(self.name), normalize_phone(self.phone)

    @property
    def last4(self) -> str:
        digits = re.sub(r"\D", "", self.phone)
        return digits[-4:] if len(digits) >= 4 else digits

    @property
    def display(self) -> str:
        return OUTPUT_NAME_FORMAT.format(name=self.name.strip(), last4=self.last4)


@dataclass
class ProductRequest:
    product: str
    count: int


@dataclass
class DrawResult:
    product: str
    winners: List[Person]


def normalize_header(s: str) -> str:
    return re.sub(r"\s+", "", str(s or "")).lower().replace("_", "")


def normalize_name(s: str) -> str:
    return re.sub(r"\s+", "", str(s or "")).strip().lower()


def normalize_phone(s: str) -> str:
    return re.sub(r"\D", "", str(s or ""))


def normalize_product(s: str) -> str:
    # PPT 줄바꿈/공백 차이를 흡수하기 위해 숫자/한글/영문 위주로 비교
    return re.sub(r"[^0-9a-zA-Z가-힣]+", "", str(s or "")).lower()


def generate_random_seed() -> int:
    """재현 가능한 추첨 로그용 10자리 랜덤 시드 생성."""
    return random.SystemRandom().randint(1000000000, 9999999999)


def format_winner_text(winners: List[Person], per_line: int = MAX_WINNERS_PER_LINE) -> str:
    """당첨자를 공백 1칸으로 구분하되, PPT에서는 한 줄에 per_line명씩 줄바꿈."""
    displays = [w.display for w in winners]
    if per_line <= 0:
        return WINNER_SEPARATOR.join(displays)
    lines = []
    for i in range(0, len(displays), per_line):
        lines.append(WINNER_SEPARATOR.join(displays[i:i + per_line]))
    return "\n".join(lines)


def read_csv_flexible(path: Path) -> List[Dict[str, str]]:
    encodings = ["utf-8-sig", "cp949", "euc-kr", "utf-8"]
    last_error = None
    for enc in encodings:
        try:
            with path.open("r", encoding=enc, newline="") as f:
                sample = f.read(4096)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=",\t;")
                except Exception:
                    dialect = csv.excel
                reader = csv.DictReader(f, dialect=dialect)
                return [{k: (v or "") for k, v in row.items()} for row in reader]
        except Exception as e:
            last_error = e
    raise RuntimeError(f"CSV 파일을 읽지 못했습니다: {path}\n{last_error}")


def find_col(row: Dict[str, str], candidates: Sequence[str]) -> Optional[str]:
    if not row:
        return None
    norm_map = {normalize_header(k): k for k in row.keys() if k is not None}
    for c in candidates:
        key = normalize_header(c)
        if key in norm_map:
            return norm_map[key]
    # 후보어가 포함된 컬럼명도 허용
    for nk, original in norm_map.items():
        for c in candidates:
            ck = normalize_header(c)
            if ck and (ck in nk or nk in ck):
                return original
    return None


def load_people(path: Path) -> List[Person]:
    rows = read_csv_flexible(path)
    if not rows:
        return []
    name_col = find_col(rows[0], NAME_HEADERS)
    phone_col = find_col(rows[0], PHONE_HEADERS)
    if not name_col or not phone_col:
        raise RuntimeError(f"{path.name}에서 이름/전화번호 컬럼을 찾지 못했습니다.\n허용 이름 컬럼: {NAME_HEADERS}\n허용 전화 컬럼: {PHONE_HEADERS}")

    people: List[Person] = []
    for idx, row in enumerate(rows, start=2):  # 헤더가 1행이므로 데이터 첫 줄은 2행
        name = (row.get(name_col) or "").strip()
        phone = (row.get(phone_col) or "").strip()
        if not name or not normalize_phone(phone):
            continue
        people.append(Person(source_row_no=idx, name=name, phone=phone, raw_row=dict(row)))
    return people


def load_products(path: Path) -> Tuple[List[ProductRequest], Optional[bool]]:
    rows = read_csv_flexible(path)
    if not rows:
        raise RuntimeError("상품.csv가 비어 있습니다.")
    product_col = find_col(rows[0], PRODUCT_HEADERS)
    count_col = find_col(rows[0], COUNT_HEADERS)
    dup_col = find_col(rows[0], DUP_HEADERS)
    if not product_col or not count_col:
        raise RuntimeError(f"상품.csv에서 상품명/인원수 컬럼을 찾지 못했습니다.\n예: 상품명,인원수,중복허용")

    products: List[ProductRequest] = []
    duplicate_mode: Optional[bool] = None

    for row in rows:
        product = (row.get(product_col) or "").strip()
        count_raw = (row.get(count_col) or "").strip()
        if not product:
            continue
        try:
            count = int(float(re.sub(r"[^0-9.]", "", count_raw)))
        except Exception:
            raise RuntimeError(f"상품 '{product}'의 인원수 값이 숫자가 아닙니다: {count_raw}")
        if count < 0:
            raise RuntimeError(f"상품 '{product}'의 인원수는 0 이상이어야 합니다.")
        products.append(ProductRequest(product=product, count=count))

        if dup_col:
            v = normalize_name(row.get(dup_col, ""))
            if v in YES_VALUES:
                duplicate_mode = True
            elif v in NO_VALUES:
                duplicate_mode = False

    if not products:
        raise RuntimeError("상품.csv에서 추첨할 상품이 없습니다.")
    return products, duplicate_mode


def exclude_blocked(all_people: List[Person], blocked: List[Person]) -> List[Person]:
    blocked_keys = {p.key for p in blocked}
    seen = set()
    participants = []
    for p in all_people:
        if p.key in blocked_keys:
            continue
        # 전체.csv 내부 완전 중복은 1명으로 처리
        if p.key in seen:
            continue
        seen.add(p.key)
        participants.append(p)
    return participants


def draw_winners(participants: List[Person], products: List[ProductRequest], allow_cross_product_duplicate: bool) -> List[DrawResult]:
    if not participants:
        raise RuntimeError("참여자가 0명입니다. 전체.csv와 블락.csv를 확인하세요.")

    # 요청사항 반영: 중복은 항상 불가능하게 고정합니다.
    # 상품.csv나 UI에서 중복 허용 값이 들어와도 한 사람이 여러 상품에 당첨되지 않습니다.
    allow_cross_product_duplicate = False

    results: List[DrawResult] = []
    globally_used = set()

    for req in products:
        pool = [p for p in participants if p.key not in globally_used]
        # 한 상품 안에서도 항상 중복 당첨 불가: random.sample 사용
        if req.count > len(pool):
            raise RuntimeError(
                f"'{req.product}' 추첨 인원({req.count}명)이 가능한 참여자 수({len(pool)}명)보다 많습니다.\n"
                f"현재 버전은 전체 상품 기준 중복 당첨이 불가능하므로 앞 상품 당첨자를 제외해서 부족할 수 있습니다."
            )
        winners = random.sample(pool, req.count) if req.count else []
        for w in winners:
            globally_used.add(w.key)
        results.append(DrawResult(product=req.product, winners=winners))

    return results


def iter_table_shapes(shapes):
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_table_shapes(shape.shapes)
        elif getattr(shape, "has_table", False):
            yield shape


def _copy_run_style(src_run, dst_run):
    """python-pptx가 직접 제공하지 않는 run 서식 복사를 XML 단위로 처리합니다."""
    try:
        src_rPr = src_run._r.rPr
        if src_rPr is not None:
            dst_rPr = dst_run._r.get_or_add_rPr()
            # 기존 rPr 하위요소 제거 후 원본 복사
            for child in list(dst_rPr):
                dst_rPr.remove(child)
            for child in src_rPr:
                dst_rPr.append(deepcopy(child))
            for key, value in src_rPr.attrib.items():
                dst_rPr.set(key, value)
    except Exception:
        pass


def _copy_paragraph_style(src_p, dst_p):
    """문단 정렬/줄간격/여백 같은 pPr 서식을 XML 단위로 복사합니다."""
    try:
        src_pPr = src_p._p.pPr
        if src_pPr is not None:
            dst_pPr = dst_p._p.get_or_add_pPr()
            for child in list(dst_pPr):
                dst_pPr.remove(child)
            for child in src_pPr:
                dst_pPr.append(deepcopy(child))
            for key, value in src_pPr.attrib.items():
                dst_pPr.set(key, value)
    except Exception:
        pass


def _clear_paragraph_runs(paragraph):
    for run in paragraph.runs:
        run.text = ""


def set_cell_text(cell, text: str):
    """
    당첨자_1.pptx처럼 셀 안에 들어있는 테스트(1234) 문구의 서식을 그대로 유지하면서
    실제 당첨자 텍스트만 교체합니다.

    핵심:
    - cell.text = ... 를 쓰지 않습니다. 이 방식은 PPT 서식을 깨뜨립니다.
    - 기존 문단 수만큼 각 줄을 배치합니다.
    - 부족한 문단은 첫 문단의 XML 서식을 복사해서 추가합니다.
    - 남는 기존 테스트 문구는 빈 문자열로 지웁니다.
    """
    tf = cell.text_frame
    lines = str(text or "").split("\n") if text else [""]

    if not tf.paragraphs:
        cell.text = text
        return

    template_p = tf.paragraphs[0]
    template_run = template_p.runs[0] if template_p.runs else None

    # 필요한 문단 수만큼 추가. 새 문단도 원본 문단 스타일을 최대한 복사합니다.
    while len(tf.paragraphs) < len(lines):
        new_p = tf.add_paragraph()
        _copy_paragraph_style(template_p, new_p)
        new_run = new_p.add_run()
        if template_run is not None:
            _copy_run_style(template_run, new_run)

    for idx, paragraph in enumerate(tf.paragraphs):
        if idx < len(lines):
            # 문단 안 첫 run에만 값을 넣고 나머지 run은 제거 대신 빈값 처리해서 서식 구조를 유지합니다.
            if paragraph.runs:
                paragraph.runs[0].text = lines[idx]
                if template_run is not None:
                    _copy_run_style(template_run, paragraph.runs[0])
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                run = paragraph.add_run()
                if template_run is not None:
                    _copy_run_style(template_run, run)
                run.text = lines[idx]
        else:
            _clear_paragraph_runs(paragraph)


def chunk_people(people: List[Person], chunk_size: int = MAX_WINNERS_PER_SLIDE) -> List[List[Person]]:
    if chunk_size <= 0:
        return [people]
    return [people[i:i + chunk_size] for i in range(0, len(people), chunk_size)] or [[]]


@dataclass
class ProductSlot:
    """PPT 안에서 상품명 오른쪽 당첨자 칸의 위치."""
    slide: object
    table_index: int
    row_index: int
    product_col_index: int

    @property
    def target_col_index(self) -> int:
        return self.product_col_index + 1


def get_safe_slide_layout(prs):
    """
    PPT마다 slide_layouts[6] 같은 빈 레이아웃이 없을 수 있습니다.
    어차피 새 슬라이드의 기본 placeholder를 모두 제거하므로,
    존재하는 아무 레이아웃이나 안전하게 사용합니다.
    """
    for idx in (6, 0):
        try:
            return prs.slide_layouts[idx]
        except IndexError:
            pass
    if len(prs.slide_layouts) > 0:
        return prs.slide_layouts[len(prs.slide_layouts) - 1]
    raise RuntimeError("PPT에 사용 가능한 슬라이드 레이아웃이 없습니다.")


def duplicate_slide(prs, slide):
    """
    원본 슬라이드를 최대한 동일하게 복제합니다.

    이전 방식은 슬라이드 안의 shape만 복사했기 때문에 배경, 그림, 일부 서식,
    관계 파일이 빠지거나 깨질 수 있었습니다. 이 방식은 슬라이드 XML 전체를
    복제한 뒤 이미지/차트/하이퍼링크 관계 rId를 새 슬라이드 기준으로 다시
    연결합니다.
    """
    try:
        layout = slide.slide_layout
    except Exception:
        layout = get_safe_slide_layout(prs)

    new_slide = prs.slides.add_slide(layout)

    # 새 슬라이드 기본 XML을 원본 슬라이드 XML로 교체합니다.
    # slideLayout 관계는 새 슬라이드가 이미 갖고 있으므로 유지하고,
    # 나머지 그림/차트/하이퍼링크 관계는 아래에서 다시 생성합니다.
    new_element = new_slide._element
    for child in list(new_element):
        new_element.remove(child)
    for child in slide._element:
        new_element.append(deepcopy(child))

    # 원본 슬라이드의 관계를 새 슬라이드에 복제하고 rId를 매핑합니다.
    # layout/notes는 새 슬라이드 자체 관계와 충돌할 수 있어 제외합니다.
    rid_map = {}
    for rel in slide.part.rels.values():
        reltype = rel.reltype or ""
        if "slideLayout" in reltype or "notesSlide" in reltype:
            continue
        try:
            is_external = bool(getattr(rel, "is_external", False))
            target = rel.target_ref if is_external else rel.target_part
            new_rid = new_slide.part.rels._add_relationship(rel.reltype, target, is_external)
            rid_map[rel.rId] = new_rid
        except Exception:
            pass

    if rid_map:
        for element in new_slide.element.iter():
            for attr_name, attr_value in list(element.attrib.items()):
                if attr_value in rid_map:
                    element.attrib[attr_name] = rid_map[attr_value]

    return new_slide

def get_table_by_index(slide, table_index: int):
    tables = list(iter_table_shapes(slide.shapes))
    if table_index >= len(tables):
        return None
    return tables[table_index].table


def get_target_cell_from_slot(slide, slot: ProductSlot):
    table = get_table_by_index(slide, slot.table_index)
    if table is None:
        return None
    if slot.row_index >= len(table.rows):
        return None
    row = table.rows[slot.row_index]
    if slot.target_col_index >= len(row.cells):
        return None
    return row.cells[slot.target_col_index]


def find_product_slots(slides, product_norms: set) -> Dict[str, List[ProductSlot]]:
    """
    기존 PPT에서 상품명 위치를 모두 찾습니다.
    같은 상품명이 여러 번 나오면 각각 독립 슬롯으로 보관합니다.
    예: 로보락이 두 번 있으면 첫 번째 로보락, 두 번째 로보락을 서로 다른 상품 행으로 처리합니다.
    """
    slots_by_product: Dict[str, List[ProductSlot]] = {p: [] for p in product_norms}

    for slide in slides:
        for table_index, table_shape in enumerate(iter_table_shapes(slide.shapes)):
            table = table_shape.table
            for row_index, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_text_norm = normalize_product(cell.text)
                    if not cell_text_norm:
                        continue
                    for prod_norm in product_norms:
                        if prod_norm and (
                            prod_norm == cell_text_norm
                            or prod_norm in cell_text_norm
                            or cell_text_norm in prod_norm
                        ):
                            if col_idx + 1 < len(row.cells):
                                slots_by_product.setdefault(prod_norm, []).append(
                                    ProductSlot(slide=slide, table_index=table_index, row_index=row_index, product_col_index=col_idx)
                                )
                            break
    return slots_by_product


def clear_all_result_cells_on_slide(slide, product_norms: set):
    """복제된 슬라이드에 기존 당첨자 텍스트가 따라오지 않도록 모든 상품의 오른쪽 칸을 비웁니다."""
    for table_shape in iter_table_shapes(slide.shapes):
        table = table_shape.table
        for row in table.rows:
            for col_idx, cell in enumerate(row.cells):
                cell_text_norm = normalize_product(cell.text)
                if not cell_text_norm:
                    continue
                for prod_norm in product_norms:
                    if prod_norm and (
                        prod_norm == cell_text_norm
                        or prod_norm in cell_text_norm
                        or cell_text_norm in prod_norm
                    ):
                        target_idx = col_idx + 1
                        if target_idx < len(row.cells):
                            set_cell_text(row.cells[target_idx], "")
                        break


def fill_slot(slide, slot: ProductSlot, result: DrawResult) -> bool:
    target_cell = get_target_cell_from_slot(slide, slot)
    if target_cell is None:
        return False
    winner_text = format_winner_text(result.winners)
    set_cell_text(target_cell, winner_text)
    return True


def move_slide_after_source(prs, slide_to_move, source_slide, offset_after_source: int = 0):
    """
    duplicate_slide()은 새 슬라이드를 맨 뒤에 추가합니다.
    25명 초과 페이지는 해당 상품의 원본 페이지 바로 다음에 나와야 하므로
    내부 슬라이드 ID 순서를 이동합니다.
    """
    slides = list(prs.slides)
    try:
        source_idx = slides.index(source_slide)
        move_idx = slides.index(slide_to_move)
    except ValueError:
        return

    target_idx = source_idx + 1 + offset_after_source
    sld_id_lst = prs.slides._sldIdLst
    sld_id = sld_id_lst[move_idx]
    sld_id_lst.remove(sld_id)

    if move_idx < target_idx:
        target_idx -= 1
    if target_idx < 0:
        target_idx = 0
    if target_idx > len(sld_id_lst):
        target_idx = len(sld_id_lst)
    sld_id_lst.insert(target_idx, sld_id)


def _save_and_reload_presentation(prs, temp_path: Path):
    """복제한 슬라이드의 shape/table 캐시를 갱신하기 위해 임시 저장 후 다시 엽니다."""
    temp_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(temp_path))
    return Presentation(str(temp_path))


def ensure_slots_for_results(prs, results: List[DrawResult], product_norms: set) -> None:
    """
    각 상품의 당첨자 수에 맞춰 필요한 슬롯 수를 먼저 확보합니다.

    중요:
    - 기존 코드는 1~25번을 먼저 입력한 슬라이드를 복사해서 복사본에도 같은 사람이 들어갔습니다.
    - 여기서는 어떤 당첨자도 입력하기 전에, 원본 템플릿 상태의 슬라이드를 먼저 복사합니다.
    - 복사 후에는 저장/재로딩 단계에서 표/도형 캐시를 갱신한 뒤 실제 값을 입력합니다.
    """
    slots_by_product = find_product_slots(list(prs.slides), product_norms)
    used_slot_count: Dict[str, int] = {prod_norm: 0 for prod_norm in product_norms}
    inserted_after_count: Dict[int, int] = {}

    for result in results:
        prod_norm = normalize_product(result.product)
        chunks = chunk_people(result.winners, MAX_WINNERS_PER_SLIDE) or [[]]
        needed = len(chunks)

        available_slots = slots_by_product.get(prod_norm, [])
        slot_idx = used_slot_count.get(prod_norm, 0)
        remaining = len(available_slots) - slot_idx

        if remaining <= 0:
            # PPT 안에서 상품명 자체를 찾지 못한 경우. 실제 unmatched 처리는 fill 단계에서 합니다.
            continue

        shortage = needed - remaining
        if shortage > 0:
            # 이번 상품이 실제로 사용할 첫 슬롯 또는 마지막 슬롯을 복제 기준으로 사용합니다.
            # 이 시점에는 아직 당첨자를 입력하지 않았으므로 복사본에 1~25번이 따라가지 않습니다.
            source_slot = available_slots[min(slot_idx, len(available_slots) - 1)]
            source_slide = source_slot.slide
            source_key = id(source_slide)
            inserted_after_count.setdefault(source_key, 0)

            for _ in range(shortage):
                new_slide = duplicate_slide(prs, source_slide)
                move_slide_after_source(prs, new_slide, source_slide, inserted_after_count[source_key])
                inserted_after_count[source_key] += 1

            # 새로 추가된 슬롯까지 포함해 슬롯 목록을 갱신합니다.
            slots_by_product = find_product_slots(list(prs.slides), product_norms)
            available_slots = slots_by_product.get(prod_norm, [])

        used_slot_count[prod_norm] = slot_idx + needed


def fill_ppt(template_path: Path, output_path: Path, results: List[DrawResult]) -> List[str]:
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다. pip install python-pptx 를 실행하세요.")

    prs = Presentation(str(template_path))
    product_norms = {normalize_product(r.product) for r in results}

    # 1) 먼저 필요한 페이지/슬롯을 모두 확보합니다. 아직 당첨자 텍스트는 입력하지 않습니다.
    ensure_slots_for_results(prs, results, product_norms)

    # 2) 복제된 슬라이드의 표/도형 캐시를 갱신합니다.
    output_path.parent.mkdir(parents=True, exist_ok=True)
    temp_path = output_path.parent / f".__tmp_{output_path.stem}.pptx"
    prs = _save_and_reload_presentation(prs, temp_path)

    # 3) 갱신된 전체 슬라이드에서 상품 슬롯을 다시 찾고, chunk 순서대로 정확히 입력합니다.
    slots_by_product = find_product_slots(list(prs.slides), product_norms)
    used_slot_count: Dict[str, int] = {prod_norm: 0 for prod_norm in product_norms}
    unmatched: List[str] = []

    for result in results:
        prod_norm = normalize_product(result.product)
        chunks = chunk_people(result.winners, MAX_WINNERS_PER_SLIDE) or [[]]
        available_slots = slots_by_product.get(prod_norm, [])
        slot_idx = used_slot_count.get(prod_norm, 0)

        if slot_idx >= len(available_slots):
            unmatched.append(result.product)
            continue

        for chunk_index, winners_chunk in enumerate(chunks):
            target_index = slot_idx + chunk_index
            if target_index >= len(available_slots):
                unmatched.append(result.product)
                break
            slot = available_slots[target_index]
            fill_slot(slot.slide, slot, DrawResult(product=result.product, winners=winners_chunk))

        used_slot_count[prod_norm] = slot_idx + len(chunks)

    # 4) 템플릿에 남아 있는 미사용 슬롯의 테스트 문구는 빈칸 처리합니다.
    for prod_norm, slots in slots_by_product.items():
        used_until = used_slot_count.get(prod_norm, 0)
        for slot in slots[used_until:]:
            target_cell = get_target_cell_from_slot(slot.slide, slot)
            if target_cell is not None:
                set_cell_text(target_cell, "")

    prs.save(str(output_path))
    try:
        temp_path.unlink(missing_ok=True)
    except Exception:
        pass
    return unmatched

def next_output_index(base_folder: Path) -> int:
    """
    당첨자/당첨자_N.pptx와 블락/블락_N.csv를 함께 확인해서 다음 번호를 정합니다.
    PPT와 블락 파일 번호가 같은 회차로 맞춰지도록 하기 위함입니다.
    """
    winners_dir = base_folder / "당첨자"
    block_dir = base_folder / "블락"
    existing_nums = []
    patterns = [
        re.compile(r"^당첨자_(\d+)\.pptx$", re.IGNORECASE),
        re.compile(r"^블락_(\d+)\.csv$", re.IGNORECASE),
    ]
    for folder in [winners_dir, block_dir]:
        if folder.exists():
            for file_path in folder.iterdir():
                for pattern in patterns:
                    m = pattern.match(file_path.name)
                    if m:
                        existing_nums.append(int(m.group(1)))
                        break
    return (max(existing_nums) + 1) if existing_nums else 1


def read_csv_with_headers(path: Path) -> Tuple[List[str], List[Dict[str, str]]]:
    """CSV를 읽을 때 헤더 순서를 함께 보존합니다."""
    encodings = ["utf-8-sig", "cp949", "euc-kr", "utf-8"]
    last_error = None
    for enc in encodings:
        try:
            with path.open("r", encoding=enc, newline="") as f:
                sample = f.read(4096)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=",\t;")
                except Exception:
                    dialect = csv.excel
                reader = csv.DictReader(f, dialect=dialect)
                headers = list(reader.fieldnames or [])
                rows = [{k: (v or "") for k, v in row.items()} for row in reader]
                return headers, rows
        except Exception as e:
            last_error = e
    raise RuntimeError(f"CSV 파일을 읽지 못했습니다: {path}\n{last_error}")


def get_block_headers_and_rows(block_csv_path: Path, all_people: List[Person]) -> Tuple[List[str], List[Dict[str, str]]]:
    """
    블락_N.csv는 기존 블락.csv의 내용을 유지한 채 뒤에 당첨자를 추가해야 합니다.
    따라서 헤더 순서와 기존 행을 함께 읽어옵니다.
    """
    headers: List[str] = []
    existing_rows: List[Dict[str, str]] = []

    if block_csv_path.exists():
        try:
            headers, existing_rows = read_csv_with_headers(block_csv_path)
        except Exception:
            headers, existing_rows = [], []

    if not headers and all_people:
        headers = [h for h in (all_people[0].raw_row or {}).keys() if h is not None]

    # 상품 컬럼은 한 번만 유지합니다.
    has_product_col = any(normalize_header(h) == normalize_header("상품") for h in headers)
    if not has_product_col:
        headers.append("상품")

    return headers, existing_rows


def write_block_csv(path: Path, results: List[DrawResult], block_csv_path: Path, all_people: List[Person]):
    """
    블락/블락_N.csv를 생성합니다.
    - 기존 블락.csv 내용은 그대로 먼저 보존합니다.
    - 그 아래에 이번 당첨자의 전체.csv 원본 정보 + 당첨 상품명을 추가합니다.
    - 헤더는 기존 블락.csv 헤더를 기준으로 하고, 없으면 전체.csv 헤더를 사용합니다.
    """
    headers, existing_rows = get_block_headers_and_rows(block_csv_path, all_people)
    path.parent.mkdir(parents=True, exist_ok=True)

    product_header = None
    for h in headers:
        if normalize_header(h) == normalize_header("상품"):
            product_header = h
            break
    if product_header is None:
        product_header = "상품"
        headers.append(product_header)

    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=headers, extrasaction="ignore")
        writer.writeheader()

        # 1) 기존 블락.csv 내용을 먼저 그대로 출력합니다.
        for row in existing_rows:
            out = {h: row.get(h, "") for h in headers}
            writer.writerow(out)

        # 2) 이번 당첨자를 이어서 출력합니다.
        for result in results:
            for winner in result.winners:
                raw = winner.raw_row or {}
                out = {h: raw.get(h, "") for h in headers}
                out[product_header] = result.product
                writer.writerow(out)

def run_job(base_folder: Path, allow_duplicate_from_ui: bool, seed: Optional[int] = None) -> str:
    if seed is None:
        seed = generate_random_seed()
    random.seed(seed)

    all_csv = base_folder / "전체.csv"
    block_csv = base_folder / "블락.csv"
    product_csv = base_folder / "상품.csv"
    # 기본 파일명은 당첨자.pptx입니다.
    # 다만 같은 폴더에 당첨자_1.pptx가 있으면, 사용자가 요청한 기준 서식으로 보고 우선 사용합니다.
    # 배포/운영 시에는 기준 템플릿 파일명을 당첨자.pptx로 바꿔 두는 것을 권장합니다.
    template_pptx = base_folder / "당첨자.pptx"
    preferred_template_pptx = base_folder / "당첨자_1.pptx"
    if preferred_template_pptx.exists():
        template_pptx = preferred_template_pptx

    missing = [p.name for p in [all_csv, block_csv, product_csv, template_pptx] if not p.exists()]
    if missing:
        raise RuntimeError("필수 파일이 없습니다: " + ", ".join(missing))

    all_people = load_people(all_csv)
    blocked_people = load_people(block_csv) if block_csv.exists() else []
    participants = exclude_blocked(all_people, blocked_people)
    products, duplicate_from_csv = load_products(product_csv)
    # 요청사항 반영: 중복은 항상 불가능하게 고정합니다.
    allow_duplicate = False

    results = draw_winners(participants, products, allow_duplicate)

    idx = next_output_index(base_folder)
    output_pptx = base_folder / "당첨자" / f"당첨자_{idx}.pptx"
    output_block_csv = base_folder / "블락" / f"블락_{idx}.csv"

    unmatched = fill_ppt(template_pptx, output_pptx, results)
    write_block_csv(output_block_csv, results, block_csv, all_people)

    lines = [
        "생성 완료",
        f"참여자 수: {len(participants)}명 / 전체: {len(all_people)}명 / 블락: {len(blocked_people)}명",
        f"중복 허용: 아니오(항상 중복 불가)",
        f"랜덤 시드: {seed}",
        f"PPTX: {output_pptx}",
        f"블락 CSV: {output_block_csv}",
    ]
    if unmatched:
        lines.append("주의: PPT에서 상품명을 찾지 못해 입력하지 못한 상품: " + ", ".join(unmatched))
    return "\n".join(lines)


class App:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title("당첨자 PPT 생성기")
        self.root.geometry("720x420")
        self.folder_var = tk.StringVar(value=str(Path.cwd()))
        self.dup_var = tk.BooleanVar(value=False)
        self.seed_var = tk.StringVar(value=str(generate_random_seed()))

        frm = ttk.Frame(self.root, padding=18)
        frm.pack(fill="both", expand=True)

        ttk.Label(frm, text="같은 폴더에 전체.csv / 블락.csv / 상품.csv / 당첨자.pptx를 넣어주세요.").pack(anchor="w")

        row = ttk.Frame(frm)
        row.pack(fill="x", pady=(16, 8))
        ttk.Label(row, text="작업 폴더").pack(side="left")
        ttk.Entry(row, textvariable=self.folder_var).pack(side="left", fill="x", expand=True, padx=8)
        ttk.Button(row, text="폴더 선택", command=self.choose_folder).pack(side="left")

        opt = ttk.LabelFrame(frm, text="중복 당첨 설정")
        opt.pack(fill="x", pady=8)
        ttk.Label(opt, text="현재 버전은 전체 상품 기준 중복 당첨이 항상 불가능합니다.").pack(anchor="w", padx=10, pady=8)
        ttk.Label(opt, text="※ 상품.csv의 중복허용 컬럼이나 화면 선택값은 무시됩니다.").pack(anchor="w", padx=10, pady=(0, 8))

        seed_row = ttk.Frame(frm)
        seed_row.pack(fill="x", pady=8)
        ttk.Label(seed_row, text="랜덤 시드(선택)").pack(side="left")
        ttk.Entry(seed_row, textvariable=self.seed_var, width=20).pack(side="left", padx=8)
        ttk.Label(seed_row, text="10자리 랜덤값이 자동 입력됩니다. 같은 시드를 쓰면 같은 결과가 나옵니다.").pack(side="left")

        ttk.Button(frm, text="당첨자 생성", command=self.generate).pack(fill="x", pady=12)

        self.log = tk.Text(frm, height=10)
        self.log.pack(fill="both", expand=True)

    def choose_folder(self):
        folder = filedialog.askdirectory(initialdir=self.folder_var.get() or str(Path.cwd()))
        if folder:
            self.folder_var.set(folder)

    def generate(self):
        try:
            seed_text = self.seed_var.get().strip()
            seed = int(seed_text) if seed_text else None
            msg = run_job(Path(self.folder_var.get()), self.dup_var.get(), seed=seed)
            self.log.delete("1.0", "end")
            self.log.insert("end", msg)
            # 완료 알림창은 띄우지 않습니다. 로그에만 출력합니다.
        except Exception as e:
            err = str(e) + "\n\n" + traceback.format_exc()
            self.log.delete("1.0", "end")
            self.log.insert("end", err)
            # 오류 알림창은 띄우지 않습니다. 로그에만 출력합니다.

    def run(self):
        self.root.mainloop()


def get_default_base_folder() -> Path:
    # exe로 빌드된 경우 exe가 있는 폴더를 기본 작업 폴더로 사용합니다.
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path.cwd()


def has_required_files(base_folder: Path) -> bool:
    required = ["전체.csv", "블락.csv", "상품.csv", "당첨자.pptx"]
    return all((base_folder / name).exists() for name in required)


def main():
    # --cli "C:\path\to\folder" 형태로도 실행 가능
    if "--cli" in sys.argv:
        idx = sys.argv.index("--cli")
        base = Path(sys.argv[idx + 1]) if len(sys.argv) > idx + 1 else get_default_base_folder()
        seed = None
        if "--seed" in sys.argv:
            seed_idx = sys.argv.index("--seed")
            if len(sys.argv) > seed_idx + 1:
                seed = int(sys.argv[seed_idx + 1])
        print(run_job(base, False, seed=seed))
        return

    base = get_default_base_folder()

    # exe와 같은 폴더에 필수 파일 4개가 있으면 실행 즉시 결과를 생성합니다.
    # 필수 파일이 없으면 기존처럼 폴더 선택 GUI를 띄웁니다.
    if has_required_files(base):
        try:
            msg = run_job(base, False)
            print(msg)
            if tk is not None:
                root = tk.Tk()
                root.withdraw()
                # 완료 알림창은 띄우지 않습니다. 로그에만 출력합니다.
                root.destroy()
        except Exception as e:
            err = str(e) + "\n\n" + traceback.format_exc()
            print(err)
            if tk is not None:
                root = tk.Tk()
                root.withdraw()
                # 오류 알림창은 띄우지 않습니다. 로그에만 출력합니다.
                root.destroy()
        return

    if tk is None:
        print(run_job(base, allow_duplicate_from_ui=False))
    else:
        App().run()


if __name__ == "__main__":
    main()
